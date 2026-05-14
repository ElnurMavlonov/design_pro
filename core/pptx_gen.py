from __future__ import annotations

import re
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour themes ────────────────────────────────────────────────────────────
THEMES: dict[str, dict] = {
    "education": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x1A, 0x73, 0xE8),   # Google-blue
        "title_text": RGBColor(0x1A, 0x73, 0xE8),
        "body_text": RGBColor(0x20, 0x20, 0x20),
        "bar": RGBColor(0x1A, 0x73, 0xE8),
    },
    "science": {
        "bg": RGBColor(0xF0, 0xF7, 0xF4),
        "accent": RGBColor(0x0D, 0x74, 0x48),   # emerald
        "title_text": RGBColor(0x0D, 0x74, 0x48),
        "body_text": RGBColor(0x1A, 0x1A, 0x1A),
        "bar": RGBColor(0x0D, 0x74, 0x48),
    },
    "history": {
        "bg": RGBColor(0xFD, 0xF6, 0xEC),
        "accent": RGBColor(0x8B, 0x4B, 0x1A),   # warm brown
        "title_text": RGBColor(0x8B, 0x4B, 0x1A),
        "body_text": RGBColor(0x2C, 0x1A, 0x0A),
        "bar": RGBColor(0x8B, 0x4B, 0x1A),
    },
    "technology": {
        "bg": RGBColor(0xF5, 0xF5, 0xF7),
        "accent": RGBColor(0x5E, 0x35, 0xB1),   # deep purple
        "title_text": RGBColor(0x5E, 0x35, 0xB1),
        "body_text": RGBColor(0x1A, 0x1A, 0x2E),
        "bar": RGBColor(0x5E, 0x35, 0xB1),
    },
    "default": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x6C, 0x63, 0xFF),   # purple
        "title_text": RGBColor(0x6C, 0x63, 0xFF),
        "body_text": RGBColor(0x20, 0x20, 0x20),
        "bar": RGBColor(0x6C, 0x63, 0xFF),
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int, bold: bool, color: RGBColor,
                  align=PP_ALIGN.LEFT, wrap: bool = True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


# ── Slide builders ───────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, slide_data: dict, theme: dict) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg"])

    # Accent bar on left
    _add_rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, theme["bar"])

    # Decorative bottom strip
    _add_rect(slide, Inches(0.25), Inches(6.8), SLIDE_W - Inches(0.25), Inches(0.7), theme["accent"])

    # Main title
    title = slide_data.get("title", "")
    _add_text_box(
        slide, title,
        left=Inches(0.6), top=Inches(2.0),
        width=Inches(11.5), height=Inches(1.8),
        font_size=44, bold=True, color=theme["title_text"],
        align=PP_ALIGN.LEFT,
    )

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_text_box(
            slide, subtitle,
            left=Inches(0.6), top=Inches(4.0),
            width=Inches(11.5), height=Inches(0.8),
            font_size=22, bold=False, color=theme["body_text"],
            align=PP_ALIGN.LEFT,
        )


def _build_content_slide(prs: Presentation, slide_data: dict, theme: dict) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg"])

    # Top accent bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), theme["bar"])

    # Side accent
    _add_rect(slide, Inches(0), Inches(0.08), Inches(0.2), SLIDE_H - Inches(0.08), theme["accent"])

    # Slide title
    title = slide_data.get("title", "")
    _add_text_box(
        slide, title,
        left=Inches(0.4), top=Inches(0.15),
        width=Inches(12.5), height=Inches(0.85),
        font_size=30, bold=True, color=theme["title_text"],
        align=PP_ALIGN.LEFT,
    )

    # Divider line under title
    _add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04), theme["accent"])

    # Bullet points
    bullets = slide_data.get("body", [])
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(6)

            # Bullet indicator
            run_bullet = p.add_run()
            run_bullet.text = "▪  "
            run_bullet.font.size = Pt(18)
            run_bullet.font.bold = True
            run_bullet.font.color.rgb = theme["accent"]
            run_bullet.font.name = "Calibri"

            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(20)
            run_text.font.bold = False
            run_text.font.color.rgb = theme["body_text"]
            run_text.font.name = "Calibri"


# ── Public API ───────────────────────────────────────────────────────────────

def build_pptx(data: dict) -> BytesIO:
    """Convert Claude JSON response into a .pptx BytesIO buffer."""
    theme_name = data.get("theme", "default")
    theme = THEMES.get(theme_name, THEMES["default"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = data.get("slides", [])

    for idx, slide_data in enumerate(slides):
        if idx == 0:
            # First slide always rendered as title slide
            title_data = {
                "title": data.get("title", slide_data.get("title", "")),
                "subtitle": data.get("subtitle", ""),
            }
            _build_title_slide(prs, title_data, theme)
        else:
            _build_content_slide(prs, slide_data, theme)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def slugify(text: str) -> str:
    """Turn presentation title into a safe filename."""
    text = text.lower().strip()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"[\s_-]+", "_", text)
    return text[:50] or "presentation"
